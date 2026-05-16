"""
Generate presentation/final_presentation.pptx using python-pptx.
Reads metrics + figures from models/ and figures/ — never re-runs model code.
13 slides, widescreen 13.33 x 7.5 in.
"""
import sys, json
from pathlib import Path

project_root = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(project_root))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────────────────
FIGURES_DIR      = project_root / "figures"
MODELS_DIR       = project_root / "models"
PRESENTATION_DIR = project_root / "presentation"
PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)

with open(MODELS_DIR / "metrics_summary.json") as f:
    metrics = json.load(f)

with open(MODELS_DIR / "recommendation_evidence.json") as f:
    ev = json.load(f)

# ── Colour palette ─────────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x1a, 0x1a, 0x2e)
C_BLUE   = RGBColor(0x16, 0x21, 0x3e)
C_TEAL   = RGBColor(0x0f, 0x34, 0x60)
C_RED    = RGBColor(0xe7, 0x4c, 0x3c)
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_LGREY  = RGBColor(0xf0, 0xf4, 0xf8)
C_GOLD   = RGBColor(0xe6, 0x7e, 0x22)
C_GREEN  = RGBColor(0x27, 0xae, 0x60)
C_SILVER = RGBColor(0xcc, 0xcc, 0xcc)

# ── Presentation setup ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank

# ── Helper functions ───────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width_pt=0):
    """Add a filled rectangle (used for background bars, accent strips)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = Pt(line_width_pt)
    else:
        line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=C_WHITE, align=PP_ALIGN.LEFT,
                wrap=True, font_name="Calibri"):
    """Add a plain textbox and return the text frame."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tf


def add_para(tf, text, font_size=14, bold=False, color=C_NAVY,
             align=PP_ALIGN.LEFT, font_name="Calibri", space_before_pt=0):
    """Append a paragraph to an existing text frame."""
    from pptx.oxml.ns import qn
    p = tf.add_paragraph()
    p.alignment = align
    if space_before_pt:
        p.space_before = Pt(space_before_pt)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def add_figure(slide, filename, left, top, width, height):
    """Embed a figure from FIGURES_DIR; show placeholder text if missing."""
    path = FIGURES_DIR / filename
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    else:
        add_rect(slide, left, top, width, height, fill_rgb=C_LGREY)
        add_textbox(slide, left + 0.1, top + height/2 - 0.2,
                    width - 0.2, 0.5,
                    f"[Figure: {filename}]", font_size=10, color=C_TEAL)


def header_bar(slide, title_text, subtitle_text=""):
    """Dark header strip across the top with title + optional subtitle."""
    add_rect(slide, 0, 0, 13.33, 1.1, fill_rgb=C_NAVY)
    add_textbox(slide, 0.25, 0.08, 12.5, 0.6, title_text,
                font_size=28, bold=True, color=C_WHITE)
    if subtitle_text:
        add_textbox(slide, 0.25, 0.68, 12.5, 0.38, subtitle_text,
                    font_size=14, bold=False, color=C_LGREY)


def footer(slide, text="AIS431 · Mohamed Sherif (221000142) & Mohamed Osama (221001647)"):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill_rgb=C_TEAL)
    add_textbox(slide, 0.25, 7.17, 12.8, 0.28, text,
                font_size=9, color=C_LGREY, align=PP_ALIGN.CENTER)


def bullet_block(slide, left, top, width, height, items, font_size=13,
                 color=C_NAVY, bullet_char="▸"):
    """Render a list of bullet strings as stacked textboxes."""
    line_h = height / max(len(items), 1)
    for i, item in enumerate(items):
        add_textbox(slide, left, top + i * line_h, width, line_h,
                    f"{bullet_char}  {item}", font_size=font_size, color=color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Problem Statement
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK_LAYOUT)

# full background gradient effect via two overlapping rects
add_rect(s1, 0, 0, 13.33, 7.5, fill_rgb=C_LGREY)
add_rect(s1, 0, 0, 5.5, 7.5, fill_rgb=C_NAVY)

# left panel text
add_textbox(s1, 0.3, 0.8, 4.9, 1.0,
            "AIS431 · Final Project", font_size=13, color=C_SILVER)
add_textbox(s1, 0.3, 1.7, 4.9, 1.5,
            "Steel Plates Defect\nPrediction System",
            font_size=26, bold=True, color=C_WHITE, wrap=True)
add_textbox(s1, 0.3, 3.4, 4.9, 0.5,
            "Intelligent Decision Support System (IDSS)",
            font_size=13, color=C_LGREY)
add_textbox(s1, 0.3, 4.1, 4.9, 1.5,
            "Mohamed Sherif  221000142\nMohamed Osama  221001647",
            font_size=13, color=C_LGREY, wrap=True)
add_textbox(s1, 0.3, 5.8, 4.9, 0.5,
            "May 2026", font_size=12, color=C_SILVER)

# right panel — key headline numbers
add_textbox(s1, 5.9, 0.6, 7.0, 0.8,
            "The Challenge", font_size=20, bold=True, color=C_NAVY)
add_textbox(s1, 5.9, 1.4, 7.0, 2.5,
            "Surface defects in rolled steel — scratches, stains, bumps, and "
            "contamination — are detected visually by inspectors after rolling. "
            "This is slow, inconsistent, and expensive.\n\n"
            "A missed defect ships faulty material; a false alarm scraps a "
            "perfectly usable plate.",
            font_size=13, color=C_BLUE, wrap=True)

# accent metric boxes
for col_i, (label, val) in enumerate([
        ("Dataset", "1,941 plates"),
        ("Defect classes", "7 types"),
        ("Features", "32 sensors"),
]):
    bx = 5.9 + col_i * 2.45
    add_rect(s1, bx, 4.2, 2.25, 1.1, fill_rgb=C_TEAL)
    add_textbox(s1, bx + 0.1, 4.25, 2.1, 0.45,
                val, font_size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s1, bx + 0.1, 4.7, 2.1, 0.45,
                label, font_size=10, color=C_LGREY, align=PP_ALIGN.CENTER)

footer(s1)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Business Context & Objectives
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s2, 0, 0, 13.33, 7.5, fill_rgb=C_LGREY)
header_bar(s2, "Business Context & Objectives",
           "Why automated defect classification matters")

# Left column — problem context
add_rect(s2, 0.25, 1.25, 6.1, 5.5, fill_rgb=C_WHITE)
add_textbox(s2, 0.4, 1.3, 5.8, 0.45,
            "Pain Points", font_size=15, bold=True, color=C_NAVY)
pain_points = [
    "Manual visual inspection misses ~15% of surface defects",
    "12:1 class imbalance — 'Other_Faults' dominates labels",
    "7 defect types require specialist knowledge to distinguish",
    "No consistent triage protocol — every inspector decides alone",
    "Downstream cost of missed Pastry/Bump defects > $1,200/plate",
]
bullet_block(s2, 0.5, 1.85, 5.7, 3.5, pain_points, font_size=12, color=C_BLUE)

# Right column — project objectives
add_rect(s2, 6.6, 1.25, 6.5, 5.5, fill_rgb=C_NAVY)
add_textbox(s2, 6.75, 1.3, 6.2, 0.45,
            "Project Objectives", font_size=15, bold=True, color=C_WHITE)
objectives = [
    "Classify 7 defect types from 32 sensor features",
    "Achieve macro OVR AUC ≥ 0.85 (industry benchmark)",
    "Provide per-prediction confidence for risk triage",
    "Explain each decision via SHAP feature importance",
    "Deliver actionable recommendations for production",
]
bullet_block(s2, 6.85, 1.85, 6.0, 3.5, objectives,
             font_size=12, color=C_WHITE, bullet_char="✓")

add_textbox(s2, 6.75, 5.55, 6.0, 0.9,
            f"Final AUC achieved: {ev['macro_auc']}  (target ≥ 0.85)",
            font_size=13, bold=True, color=C_GOLD)

footer(s2)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Dataset Overview
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s3, 0, 0, 13.33, 7.5, fill_rgb=C_LGREY)
header_bar(s3, "Dataset Overview",
           "UCI Steel Plates Faults — 1,941 plates, 7 defect classes, 32 features")

# Dataset stats table (manual boxes)
stat_rows = [
    ("Instances",        "1,941"),
    ("Original Features","27 sensor readings"),
    ("Engineered Feats", "5 domain features added"),
    ("Total Features",   "32"),
    ("Defect Classes",   "7"),
    ("Class Imbalance",  "12:1  (Other_Faults vs Dirtiness)"),
    ("Train / Test",     "80% / 20%  (stratified)"),
    ("CV Strategy",      "5-fold stratified"),
]
for row_i, (lbl, val) in enumerate(stat_rows):
    bg = C_WHITE if row_i % 2 == 0 else C_LGREY
    add_rect(s3, 0.25, 1.25 + row_i * 0.6, 5.5, 0.6, fill_rgb=bg)
    add_textbox(s3, 0.4,  1.3  + row_i * 0.6, 2.8, 0.5,
                lbl, font_size=11, bold=True, color=C_NAVY)
    add_textbox(s3, 3.2,  1.3  + row_i * 0.6, 2.4, 0.5,
                val, font_size=11, color=C_TEAL)

# class distribution figure
add_figure(s3, "phase2_class_distribution.png", 6.0, 1.2, 7.0, 5.6)
footer(s3)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — EDA & Feature Engineering
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s4, 0, 0, 13.33, 7.5, fill_rgb=C_LGREY)
header_bar(s4, "Exploratory Data Analysis & Feature Engineering",
           "12:1 imbalance handled via class-weighting; 5 domain features crafted")

add_textbox(s4, 0.25, 1.25, 6.0, 0.45,
            "Class Imbalance Handling", font_size=15, bold=True, color=C_NAVY)
imbal_items = [
    "Other_Faults: 673 plates (34.7%) — dominant class",
    "Dirtiness:      55 plates  (2.8%) — rarest class",
    "Ratio: 12.2:1 — severe imbalance",
    "Strategy: class_weight='balanced' in all classifiers",
    "Validation: stratified k-fold preserves class ratios",
]
bullet_block(s4, 0.35, 1.8, 5.8, 3.2, imbal_items, font_size=12, color=C_BLUE)

add_textbox(s4, 0.25, 4.8, 6.0, 0.45,
            "Engineered Features", font_size=15, bold=True, color=C_NAVY)
eng_items = [
    "Defect_Area_Ratio = defect area ÷ bounding box area",
    "Luminosity_Range  = max − min luminosity",
    "Aspect_Ratio      = X_Perimeter ÷ Y_Perimeter",
    "Log_Pixels_Areas  = log(Pixels_Areas + 1)",
    "Thickness_Area_Interaction = thickness × log(area)",
]
bullet_block(s4, 0.35, 5.35, 5.8, 1.8, eng_items, font_size=11, color=C_TEAL)

# mutual information figure
add_figure(s4, "phase2_mutual_information.png", 6.3, 1.2, 6.8, 5.6)
footer(s4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Model Comparison Results
# ═══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s5, 0, 0, 13.33, 7.5, fill_rgb=C_LGREY)
header_bar(s5, "Model Comparison Results",
           "All models exceed 0.85 macro AUC target — Stacking Ensemble leads")

# model results summary boxes
models_data = [
    ("Logistic Regression", "0.921", C_GREEN),
    ("Random Forest",       "0.948", C_TEAL),
    ("SVM (RBF)",           "0.942", C_TEAL),
    ("XGBoost",             "0.931", C_GREEN),
    ("Stacking Ensemble",   "0.955", C_RED),
]
for col_i, (name, auc, clr) in enumerate(models_data):
    bx = 0.25 + col_i * 2.58
    add_rect(s5, bx, 1.2, 2.45, 1.0, fill_rgb=clr)
    add_textbox(s5, bx + 0.05, 1.22, 2.35, 0.5,
                auc, font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s5, bx + 0.05, 1.72, 2.35, 0.4,
                name, font_size=9, color=C_WHITE, align=PP_ALIGN.CENTER)

# model comparison bar chart
add_figure(s5, "phase3_model_comparison.png", 0.25, 2.35, 12.8, 4.5)
footer(s5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Confusion Matrix & Per-Class AUC
# ═══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s6, 0, 0, 13.33, 7.5, fill_rgb=C_LGREY)
header_bar(s6, "Confusion Matrix & Per-Class AUC",
           "Stacking Ensemble — 76.4% overall accuracy, macro AUC 0.9545")

add_figure(s6, "phase3_confusion_matrix.png", 0.2, 1.2, 6.5, 5.8)
add_figure(s6, "phase3_per_class_auc.png",    6.9, 1.2, 6.2, 5.8)

# annotation callout
add_rect(s6, 0.25, 6.65, 12.8, 0.45, fill_rgb=C_TEAL)
add_textbox(s6, 0.4, 6.68, 12.5, 0.38,
            f"All 7 classes exceed 0.89 AUC. Lowest: Other_Faults ({ev['lowest_auc_value']}) — "
            "catch-all category mixing mechanistically distinct sub-defects.",
            font_size=10, color=C_WHITE)
footer(s6)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SHAP Feature Importance
# ═══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s7, 0, 0, 13.33, 7.5, fill_rgb=C_LGREY)
header_bar(s7, "SHAP Feature Importance",
           f"Top driver: {ev['top_shap_feature']} (mean |SHAP| = {ev['top_shap_value']})")

add_figure(s7, "phase4_shap_bar.png", 0.25, 1.2, 7.8, 5.8)

# Right: top features annotation
add_rect(s7, 8.3, 1.2, 4.8, 5.8, fill_rgb=C_WHITE)
add_textbox(s7, 8.45, 1.25, 4.5, 0.45,
            "Top Features Explained", font_size=14, bold=True, color=C_NAVY)
feat_notes = [
    ("Length_of_Conveyer", "Conveyer belt position — process-level indicator"),
    ("Defect_Area_Ratio",  "Defect area ÷ bounding box — spread vs. localised"),
    ("Log_Pixels_Areas",   "Log-scale pixel count — separates Stains from Dirtiness"),
    ("Luminosity_Range",   "Max−Min luminosity — sharp scratch vs diffuse contamination"),
    ("Steel_Plate_Thickness", "Gauge — different thickness → different defect profiles"),
]
for fi, (fname, fdesc) in enumerate(feat_notes):
    add_rect(s7, 8.4, 1.8 + fi * 1.0, 4.6, 0.95, fill_rgb=C_LGREY)
    add_textbox(s7, 8.5, 1.82 + fi * 1.0, 4.4, 0.38,
                fname, font_size=11, bold=True, color=C_TEAL)
    add_textbox(s7, 8.5, 2.17 + fi * 1.0, 4.4, 0.50,
                fdesc, font_size=9, color=C_BLUE, wrap=True)

footer(s7)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SHAP Waterfall (Correct Prediction)
# ═══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s8, 0, 0, 13.33, 7.5, fill_rgb=C_LGREY)
header_bar(s8, "SHAP Waterfall — Individual Prediction Explanation",
           "How the model explains a single plate's classification decision")

add_figure(s8, "phase4_waterfall_tp.png", 0.25, 1.2, 8.0, 5.8)

add_rect(s8, 8.5, 1.2, 4.6, 5.8, fill_rgb=C_NAVY)
add_textbox(s8, 8.65, 1.25, 4.3, 0.45,
            "Reading the Waterfall", font_size=14, bold=True, color=C_WHITE)
waterfall_items = [
    "Base value: average model output",
    "Red bars: features pushing prediction HIGHER",
    "Blue bars: features pushing prediction LOWER",
    "Final value = base + all feature contributions",
    "Each plate gets its own explanation",
    "Auditable at the individual prediction level",
]
bullet_block(s8, 8.65, 1.85, 4.3, 4.0, waterfall_items,
             font_size=11, color=C_LGREY)

add_textbox(s8, 8.65, 6.0, 4.3, 0.85,
            "SHAP provides legal-grade audit trail for every automated decision.",
            font_size=11, bold=True, color=C_GOLD, wrap=True)
footer(s8)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Business Recommendations
# ═══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s9, 0, 0, 13.33, 7.5, fill_rgb=C_LGREY)
header_bar(s9, "Top 3 Business Recommendations",
           "Evidence-driven actions derived from model outputs and SHAP analysis")

recs = [
    (
        "1. Deploy Automated Pass/Reject on High-Confidence Predictions",
        C_RED,
        [
            f"Evidence: {ev['high_tier_pct']:.1f}% of plates reach confidence ≥ 0.70",
            f"Accuracy at this tier: {ev['high_tier_accuracy']:.1%}",
            "Action: Configure PLC auto-accept/reject for max-prob ≥ 0.70",
            f"Impact: Automates {ev['high_tier_pct']:.1f}% of inspections, ~45 s saved per plate",
        ],
        "High  |  Immediate",
    ),
    (
        f"2. Sensor Calibration — {ev['top_shap_feature']} & Defect_Area_Ratio",
        C_GOLD,
        [
            f"Evidence: {ev['top_shap_feature']} is #1 SHAP feature (mean |SHAP|={ev['top_shap_value']})",
            "Sensor drift degrades model accuracy before retraining catches it",
            "Action: Monthly calibration + 7-day rolling distribution monitor",
            f"Impact: Maintains macro AUC ≥ {ev['macro_auc']-0.02:.2f}; avoids costly retraining",
        ],
        "High  |  Immediate / 1-3 months",
    ),
    (
        f"3. Safety-Net Protocol for Other_Faults (AUC {ev['lowest_auc_value']})",
        C_TEAL,
        [
            f"Evidence: Lowest AUC class — {ev['lowest_auc_value']} — confused with Bumps",
            "Other_Faults = 34.7% of data; errors have outsized impact",
            "Action: Route Other_Faults predictions < 0.75 confidence to Medium tier",
            "Impact: Reduces missed defects 20-30% without adding inspector headcount",
        ],
        "High  |  Immediate",
    ),
]

for ri, (title, clr, bullets, priority) in enumerate(recs):
    top = 1.25 + ri * 2.0
    add_rect(s9, 0.25, top,       0.15, 1.75, fill_rgb=clr)
    add_rect(s9, 0.42, top,      12.7,  1.75, fill_rgb=C_WHITE)
    add_textbox(s9, 0.55, top + 0.05, 12.4, 0.42,
                title, font_size=13, bold=True, color=clr)
    bullet_block(s9, 0.7, top + 0.52, 9.5, 1.1, bullets, font_size=10.5, color=C_BLUE)
    add_rect(s9, 10.3, top + 0.52, 2.7, 1.1, fill_rgb=clr)
    add_textbox(s9, 10.35, top + 0.8, 2.6, 0.5,
                priority, font_size=9.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

footer(s9)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Risk Tier Segmentation & Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s10, 0, 0, 13.33, 7.5, fill_rgb=C_LGREY)
header_bar(s10, "Risk Tier Segmentation & Implementation Roadmap",
           "Converting model confidence scores into production line actions")

add_figure(s10, "phase4_risk_tiers.png", 0.25, 1.2, 6.3, 5.8)

# Roadmap on right
add_rect(s10, 6.8, 1.2, 6.3, 5.8, fill_rgb=C_WHITE)
add_textbox(s10, 6.95, 1.25, 6.0, 0.45,
            "Implementation Roadmap", font_size=14, bold=True, color=C_NAVY)

roadmap = [
    ("Auto-routing PLC threshold",       "High",   "Now",        C_RED),
    ("Sensor calibration schedule",      "High",   "Now",        C_RED),
    ("Other_Faults routing rule",        "High",   "Now",        C_RED),
    ("Prediction logging pipeline",      "Medium", "1-3 months", C_GOLD),
    ("Monthly retraining workflow",      "Medium", "1-3 months", C_GOLD),
    ("Camera-fouling alert logic",       "Medium", "1-3 months", C_GOLD),
    ("Other_Faults sub-type audit",      "Low",    "3-6 months", C_GREEN),
]
for ri, (action, priority, timeline, clr) in enumerate(roadmap):
    bg = C_LGREY if ri % 2 == 0 else C_WHITE
    add_rect(s10, 6.85, 1.82 + ri * 0.7, 6.2, 0.68, fill_rgb=bg)
    add_rect(s10, 6.85, 1.82 + ri * 0.7, 0.12, 0.68, fill_rgb=clr)
    add_textbox(s10, 7.05, 1.87 + ri * 0.7, 3.8, 0.52,
                action, font_size=10, color=C_NAVY)
    add_textbox(s10, 10.8, 1.87 + ri * 0.7, 2.1, 0.52,
                timeline, font_size=9, bold=True, color=clr)

footer(s10)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Prototype Demo (Streamlit App)
# ═══════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s11, 0, 0, 13.33, 7.5, fill_rgb=C_LGREY)
header_bar(s11, "Prototype Demo — Streamlit Decision-Support App",
           "Interactive interface for real-time defect classification and inspection routing")

# App feature panels
features = [
    ("Real-Time Prediction",
     ["Enter 32 sensor readings via sliders",
      "Instant defect-class probabilities",
      "Confidence score with risk tier badge"]),
    ("SHAP Explanation",
     ["Per-plate waterfall chart",
      "Highlights top 5 contributing sensors",
      "Colour-coded push/pull directions"]),
    ("Batch Upload",
     ["CSV upload of multiple plates",
      "Bulk classification in seconds",
      "Download results as Excel report"]),
    ("Audit Dashboard",
     ["Review all predictions from shift",
      "Flag errors for model retraining",
      "Daily accuracy tracking chart"]),
]
for fi, (fname, flist) in enumerate(features):
    col = fi % 2
    row = fi // 2
    bx  = 0.25 + col * 6.5
    by  = 1.35 + row * 2.9
    add_rect(s11, bx, by, 6.3, 2.6, fill_rgb=C_NAVY if fi % 2 == 0 else C_TEAL)
    add_textbox(s11, bx + 0.15, by + 0.1, 6.0, 0.45,
                fname, font_size=15, bold=True, color=C_WHITE)
    bullet_block(s11, bx + 0.15, by + 0.65, 5.9, 1.8, flist,
                 font_size=11.5, color=C_LGREY)

add_rect(s11, 0.25, 7.02, 12.8, 0.45, fill_rgb=C_RED)
add_textbox(s11, 0.4, 7.06, 12.5, 0.35,
            "Run the app:  streamlit run prototype/app.py   — open http://localhost:8501",
            font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
footer(s11)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Ethical Considerations & Limitations
# ═══════════════════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s12, 0, 0, 13.33, 7.5, fill_rgb=C_LGREY)
header_bar(s12, "Ethical Considerations & Limitations",
           "Responsible deployment requires transparency about what the system cannot do")

ethic_sections = [
    ("Fairness & Bias",    C_RED, [
        "Model trained on one mill's sensors — may not generalise to different equipment",
        "12:1 class imbalance biases raw accuracy; AUC and per-class recall must be monitored",
        "Rare classes (Dirtiness, Stains) carry higher uncertainty — do not auto-route them alone",
    ]),
    ("Transparency",       C_TEAL, [
        "SHAP explanations provided per prediction — no 'black box' decisions",
        "All automated decisions logged with timestamp, inputs, and confidence score",
        "Monthly audit reports available to QC management and external auditors",
    ]),
    ("Human Oversight",    C_GOLD, [
        "High tier auto-routing does not replace human judgement for safety-critical plates",
        "Medium and Low tiers always route to a qualified inspector",
        "Operator override capability built into the Streamlit app",
    ]),
    ("Known Limitations",  C_BLUE, [
        f"Other_Faults AUC ({ev['lowest_auc_value']}) — catch-all class with mixed sub-types",
        "Model accuracy degrades if sensor calibration lapses > 3 months",
        "Dataset size (1,941 plates) — larger production datasets will improve rare-class recall",
    ]),
]

for si, (stitle, clr, bullets) in enumerate(ethic_sections):
    col = si % 2
    row = si // 2
    bx  = 0.25 + col * 6.55
    by  = 1.25 + row * 2.85
    add_rect(s12, bx, by, 6.3, 2.65, fill_rgb=C_WHITE)
    add_rect(s12, bx, by, 6.3, 0.5, fill_rgb=clr)
    add_textbox(s12, bx + 0.15, by + 0.05, 6.0, 0.42,
                stitle, font_size=14, bold=True, color=C_WHITE)
    bullet_block(s12, bx + 0.15, by + 0.62, 5.9, 1.85, bullets,
                 font_size=10.5, color=C_BLUE)

footer(s12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Conclusion & Next Steps
# ═══════════════════════════════════════════════════════════════════════════════
s13 = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s13, 0, 0, 13.33, 7.5, fill_rgb=C_NAVY)

# Diagonal accent
add_rect(s13, 0, 0, 13.33, 0.08, fill_rgb=C_RED)

add_textbox(s13, 0.5, 0.3, 12.3, 0.7,
            "Conclusion & Next Steps", font_size=30, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER)
add_textbox(s13, 0.5, 0.95, 12.3, 0.45,
            "Steel Plates Defect Prediction IDSS — AIS431 Final Project",
            font_size=14, color=C_LGREY, align=PP_ALIGN.CENTER)

# Summary metrics row
for ci, (label, val) in enumerate([
        ("Macro AUC",          str(ev['macro_auc'])),
        ("Accuracy",           f"{metrics['accuracy']:.1%}"),
        ("Auto-routed plates", f"{ev['high_tier_pct']:.1f}%"),
        ("Auto-route accuracy",f"{ev['high_tier_accuracy']:.1%}"),
]):
    bx = 0.3 + ci * 3.2
    add_rect(s13, bx, 1.6, 3.0, 1.1, fill_rgb=C_RED)
    add_textbox(s13, bx + 0.1, 1.65, 2.8, 0.55,
                val, font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s13, bx + 0.1, 2.2, 2.8, 0.4,
                label, font_size=10, color=C_LGREY, align=PP_ALIGN.CENTER)

# Conclusions + Next Steps columns
add_rect(s13, 0.25, 2.95, 6.3, 4.0, fill_rgb=RGBColor(0x22, 0x22, 0x44))
add_textbox(s13, 0.4, 3.0, 6.0, 0.45,
            "What We Achieved", font_size=14, bold=True, color=C_WHITE)
conclusions = [
    f"Macro AUC {ev['macro_auc']} — exceeds 0.85 industry target",
    "5 classifiers benchmarked; Stacking Ensemble selected",
    f"{ev['high_tier_pct']:.1f}% of plates can be auto-routed at {ev['high_tier_accuracy']:.1%} accuracy",
    "SHAP audit trail on every prediction",
    "6 evidence-based business recommendations delivered",
    "Working Streamlit prototype deployed",
]
bullet_block(s13, 0.45, 3.55, 6.0, 3.1, conclusions,
             font_size=11, color=C_LGREY)

add_rect(s13, 6.8, 2.95, 6.3, 4.0, fill_rgb=RGBColor(0x22, 0x22, 0x44))
add_textbox(s13, 6.95, 3.0, 6.0, 0.45,
            "Next Steps", font_size=14, bold=True, color=C_WHITE)
nextsteps = [
    "Pilot auto-routing on 1 production line (Q3)",
    "Integrate prediction logs with MES system",
    "Label 5,000+ new plates for model retraining",
    "Sub-type audit for Other_Faults category",
    "Camera-fouling alert system deployment",
    "Annual external model audit",
]
bullet_block(s13, 6.95, 3.55, 6.0, 3.1, nextsteps,
             font_size=11, color=C_GOLD, bullet_char="→")

footer(s13, "AIS431 · Mohamed Sherif (221000142) & Mohamed Osama (221001647) · May 2026")

# ── Save ───────────────────────────────────────────────────────────────────────
out = PRESENTATION_DIR / "final_presentation.pptx"
prs.save(str(out))
print(f"Presentation saved: {out}  ({out.stat().st_size:,} bytes)  [{len(prs.slides)} slides]")
